import asyncio
import logging
from datetime import datetime, timezone

from sqlalchemy import select, delete
from sqlalchemy.orm import Session

from app.database import SessionLocal
from app.models import Assignment, KnowledgeChunk
from app.ask_engine import embed_text
from app.config import settings

logger = logging.getLogger(__name__)

CHUNK_SIZE = 800      # words per chunk
CHUNK_OVERLAP = 100   # word overlap between chunks
MAX_CHUNKS_PER_FILE = 8  # cap per file so large docs don't dominate the index
INVENIAS_SUBDOMAIN = "peoplelink"

# Drives whose *content* is worth embedding for semantic search.
# Client-facing drives (Sąskaitos, Pasiūlymai, Sutartys…) are excluded —
# their files are structured/binary docs best found via search_files, not vectors.
# Override with SHAREPOINT_INDEX_DRIVES env var if needed.
_CONTENT_DRIVES_DEFAULT = [
    "Aktualūs dokumentai",
    "Akademija",
    "Mokymai",
    "Procesai ir KPI",
    "Tvarkos",
    "Instrukcijos",
    "Skelbimai",
    "Komunikacija",
    "Consulting / Tripod",
    "People Link methods",
    "Tripod methods (surveys)",
    "Assessment as a Service",
    "Klientų auginimas",
    "Kiti dokumentai",
    "Darbo taryba",
    "Darbų sauga",
    "Sveikatos patikra",
    "Asmens duomenų apsauga",
    "Laikinasis įdarbinimas",
]


def _chunk_text(text: str) -> list[str]:
    words = text.split()
    if not words:
        return []
    chunks = []
    step = CHUNK_SIZE - CHUNK_OVERLAP
    for i in range(0, len(words), step):
        chunk = " ".join(words[i : i + CHUNK_SIZE])
        if chunk.strip():
            chunks.append(chunk)
    return chunks


def _assignment_url(item_id: str) -> str:
    return f"https://{INVENIAS_SUBDOMAIN}.invenias.com/a/assignments/{item_id}"


def _assignment_text(a: Assignment) -> str:
    return (
        f"Assignment Reference: {a.reference_number}\n"
        f"Company: {a.company_name or 'Unknown'}\n"
        f"Role / Title: {a.title or 'Unknown'}\n"
        f"Status: {a.status}\n"
        f"Invenias URL: {_assignment_url(a.id)}"
    )


async def _index_assignments(db: Session) -> int:
    from app import progress as _prog
    assignments = db.execute(select(Assignment)).scalars().all()
    total = len(assignments)
    now = datetime.now(timezone.utc)
    indexed = 0

    _prog.update("indexing", phase="assignments", current_drive=f"Invenias assignments (0/{total})", drives_total=0)

    for assignment in assignments:
        content = _assignment_text(assignment)
        try:
            embedding = await embed_text(content)
        except Exception as e:
            logger.error(f"Embedding failed for assignment {assignment.id}: {e}")
            continue

        existing = db.execute(
            select(KnowledgeChunk).where(
                KnowledgeChunk.source_type == "invenias",
                KnowledgeChunk.source_id == assignment.id,
            )
        ).scalar_one_or_none()

        url = _assignment_url(assignment.id)
        if existing:
            existing.content = content
            existing.source_name = assignment.display_name
            existing.source_url = url
            existing.embedding = embedding
            existing.indexed_at = now
        else:
            db.add(KnowledgeChunk(
                source_type="invenias",
                source_id=assignment.id,
                source_name=assignment.display_name,
                source_url=url,
                content=content,
                embedding=embedding,
                indexed_at=now,
            ))
        indexed += 1
        if indexed % 10 == 0:
            _prog.update("indexing", current_drive=f"Invenias assignments ({indexed}/{total})", files_done=indexed)

    db.commit()
    return indexed


def _is_path_archived(folder_path: str, archived_paths: set[str]) -> bool:
    """Return True if folder_path equals or is under any path in the set."""
    for ap in archived_paths:
        if folder_path == ap or folder_path.startswith(ap + "/"):
            return True
    return False


def _drive_for_meta_key(item_id: str) -> str | None:
    """Resolve a cat:: or subcat:: DocMeta key to the corresponding SharePoint drive name."""
    try:
        from app.routers.documents import _CATEGORIES
    except Exception:
        return None
    if item_id.startswith("cat::"):
        cat_name = item_id[5:]
        cat = next((c for c in _CATEGORIES if c["name"] == cat_name), None)
        return cat["drive"] if cat and cat.get("drive") else None
    if item_id.startswith("subcat::"):
        parts = item_id[8:].split("::", 1)
        if len(parts) == 2:
            cat_name, sub_name = parts
            cat = next((c for c in _CATEGORIES if c["name"] == cat_name), None)
            if cat:
                sub = next((s for s in cat.get("subcategories", []) if s["name"] == sub_name), None)
                return sub["drive"] if sub else None
    return None


async def _index_one_drive(
    db: Session,
    drive_name: str,
    sp_kwargs: dict,
    http,
    now: datetime,
    changed_items: list[dict] | None = None,
    deleted_item_ids: list[str] | None = None,
    archived_item_ids: set[str] | None = None,
    archived_paths: set[str] | None = None,
    no_index_item_ids: set[str] | None = None,
    no_index_paths: set[str] | None = None,
) -> int:
    from app.sharepoint import list_files, resolve_token_and_drive

    SUPPORTED_EXTENSIONS = {".txt", ".md", ".csv", ".pdf", ".docx", ".pptx", ".xlsx"}
    MAX_FILE_BYTES = 25 * 1024 * 1024  # 25 MB — prevent OOM on Railway
    indexed = 0
    skipped_ext = 0
    skipped_size = 0

    if deleted_item_ids:
        for item_id in deleted_item_ids:
            db.execute(
                delete(KnowledgeChunk).where(
                    KnowledgeChunk.source_type == "sharepoint",
                    KnowledgeChunk.source_id.like(f"{item_id}%"),
                )
            )
        db.commit()

    if changed_items is not None and not changed_items:
        return 0

    token, drive_base = await resolve_token_and_drive(**sp_kwargs, drive_name=drive_name)
    auth_headers = {"Authorization": f"Bearer {token}"}

    if changed_items is None:
        async def collect_files(folder: str) -> list[dict]:
            try:
                items = await list_files(**sp_kwargs, drive_name=drive_name, folder=folder)
            except Exception as e:
                logger.error(f"[{drive_name}] List failed for '{folder}': {e}")
                return []
            result = []
            for item in items:
                if item["is_folder"]:
                    sub_folder = f"{folder}/{item['name']}".lstrip("/") if folder else item["name"]
                    result.extend(await collect_files(sub_folder))
                else:
                    item["folder_path"] = folder
                    result.append(item)
            return result

        all_files = await collect_files("")
        logger.info(f"[{drive_name}] Found {len(all_files)} files")
    else:
        all_files = changed_items
        logger.info(f"[{drive_name}] Processing {len(all_files)} changed files (delta)")

    import io as _io

    for f in all_files:
        name: str = f["name"]
        ext = "." + name.rsplit(".", 1)[-1].lower() if "." in name else ""
        if ext not in SUPPORTED_EXTENSIONS:
            skipped_ext += 1
            continue

        file_size = f.get("size") or 0
        if file_size > MAX_FILE_BYTES:
            logger.warning(f"[{drive_name}] Skipping {name} — {file_size:,} bytes (>{MAX_FILE_BYTES//1024//1024} MB)")
            skipped_size += 1
            continue

        # Determine archive status: file-level or folder-level
        file_folder_path = f.get("folder_path", "")
        file_is_archive = bool(
            (archived_item_ids and f["id"] in archived_item_ids) or
            (archived_paths and _is_path_archived(file_folder_path, archived_paths))
        )

        # Skip files excluded from indexing; remove any stale chunks for them
        file_is_no_index = bool(
            (no_index_item_ids and f["id"] in no_index_item_ids) or
            (no_index_paths and _is_path_archived(file_folder_path, no_index_paths))
        )
        if file_is_no_index:
            db.execute(
                delete(KnowledgeChunk).where(
                    KnowledgeChunk.source_type == "sharepoint",
                    KnowledgeChunk.source_id.like(f"{f['id']}%"),
                )
            )
            db.commit()
            continue

        # Incremental: skip if already indexed and not modified since
        file_modified: datetime | None = None
        file_modified_str = f.get("modified", "")
        if file_modified_str:
            file_modified = datetime.fromisoformat(file_modified_str.replace("Z", "+00:00"))
            existing = db.execute(
                select(KnowledgeChunk.indexed_at).where(
                    KnowledgeChunk.source_type == "sharepoint",
                    KnowledgeChunk.source_id == f["id"],
                )
            ).scalar_one_or_none()
            if existing is not None and existing.replace(tzinfo=timezone.utc) >= file_modified:
                continue  # Already up to date

        try:
            resp = await http.get(
                f"{drive_base}/items/{f['id']}/content",
                headers=auth_headers,
                follow_redirects=True,
            )
            resp.raise_for_status()
            raw = resp.content
        except Exception as e:
            logger.error(f"[{drive_name}] Download failed for {name}: {e}")
            continue

        try:
            if ext == ".pdf":
                from pypdf import PdfReader
                reader = PdfReader(_io.BytesIO(raw))
                full_text = "\n".join(p.extract_text() or "" for p in reader.pages)
            elif ext == ".docx":
                from docx import Document
                doc = Document(_io.BytesIO(raw))
                full_text = "\n".join(p.text for p in doc.paragraphs)
            elif ext == ".pptx":
                from pptx import Presentation
                prs = Presentation(_io.BytesIO(raw))
                slide_texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text)
                full_text = "\n".join(slide_texts)
            elif ext == ".xlsx":
                import openpyxl
                wb = openpyxl.load_workbook(_io.BytesIO(raw), read_only=True, data_only=True)
                rows = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        row_text = " ".join(str(c) for c in row if c is not None)
                        if row_text.strip():
                            rows.append(row_text)
                full_text = "\n".join(rows)
            else:
                full_text = raw.decode("utf-8", errors="replace")
        except Exception as e:
            logger.error(f"[{drive_name}] Text extraction failed for {name}: {e}")
            continue

        full_text = full_text.replace("\x00", "")  # PostgreSQL rejects NUL characters

        # Prepend stored description so Ask PL has richer context for this file
        from app.models import DocMeta as _DocMeta
        import json as _json
        doc_meta = db.execute(
            select(_DocMeta).where(_DocMeta.item_id == f["id"])
        ).scalar_one_or_none()
        if doc_meta and doc_meta.comment:
            audience_str = ""
            try:
                aud = _json.loads(doc_meta.audience) if doc_meta.audience else []
                if aud:
                    audience_str = f" (Audience: {', '.join(aud)})"
            except Exception:
                pass
            full_text = f"Description: {doc_meta.comment}{audience_str}\n\n{full_text}"

        raw_chunks = _chunk_text(full_text)[:MAX_CHUNKS_PER_FILE]
        if not raw_chunks:
            continue

        folder_path = f.get("folder_path", "")
        full_doc_name = f"{folder_path}/{name}" if folder_path else name
        doc_prefix = f"[Drive: {drive_name}] [File: {full_doc_name}]\n\n"
        chunks = [doc_prefix + c for c in raw_chunks]

        # Delete existing chunks for this file before re-inserting
        db.execute(
            delete(KnowledgeChunk).where(
                KnowledgeChunk.source_type == "sharepoint",
                KnowledgeChunk.source_id.like(f"{f['id']}%"),
            )
        )

        file_indexed = 0
        for i, chunk_text_content in enumerate(chunks):
            try:
                embedding = await embed_text(chunk_text_content)
            except Exception as e:
                logger.error(f"[{drive_name}] Embedding failed for {name} chunk {i}: {e}")
                continue

            chunk_id = f"{f['id']}::{i}" if i > 0 else f["id"]
            db.add(KnowledgeChunk(
                source_type="sharepoint",
                source_id=chunk_id,
                source_name=full_doc_name,
                source_url=f.get("web_url"),
                content=chunk_text_content,
                embedding=embedding,
                modified=file_modified.replace(tzinfo=None) if file_modified else None,
                indexed_at=now,
                is_archive=file_is_archive,
            ))
            file_indexed += 1

        # Commit per file and clear identity map to prevent memory accumulation
        db.commit()
        db.expunge_all()
        indexed += file_indexed

    logger.info(f"[{drive_name}] Skipped {skipped_ext} unsupported + {skipped_size} oversized files")
    return indexed


def _all_category_drives() -> list[str]:
    from app.routers.documents import _all_category_drives as _cat_drives
    return _cat_drives()


async def _collect_drive_files(sp_kwargs: dict, drive_name: str, folder: str = "") -> list[dict]:
    """Recursively list every file in a drive (listing only — no downloads)."""
    from app.sharepoint import list_files

    try:
        items = await list_files(**sp_kwargs, drive_name=drive_name, folder=folder)
    except Exception as e:
        logger.error(f"[catalog:{drive_name}] List failed for '{folder}': {e}")
        return []

    result: list[dict] = []
    for item in items:
        if item.get("is_folder"):
            sub = f"{folder}/{item['name']}".lstrip("/") if folder else item["name"]
            result.extend(await _collect_drive_files(sp_kwargs, drive_name, sub))
        else:
            item["folder_path"] = folder
            result.append(item)
    return result


async def _build_file_catalog(db: Session) -> int:
    """Catalog every file across all category drives for fast filename/company search.

    Cheap and memory-light: it only lists files, never downloads or embeds them.
    Runs before content indexing so file search works even if embedding is slow.
    """
    if not all([
        settings.sharepoint_tenant_id,
        settings.sharepoint_client_id,
        settings.sharepoint_client_secret,
    ]):
        return 0

    from app.sharepoint import _normalize, list_files_delta, DeltaExpiredError
    from app.models import FileCatalog, DriveSync
    from app import progress as _prog

    # Catalog always covers every drive so search_files can find any document
    drives = _all_category_drives()

    sp_kwargs = dict(
        tenant_id=settings.sharepoint_tenant_id,
        client_id=settings.sharepoint_client_id,
        client_secret=settings.sharepoint_client_secret,
        site_hostname=settings.sharepoint_site_hostname,
        site_path=settings.sharepoint_site_path,
    )

    now = datetime.now(timezone.utc)
    total = 0

    _prog.update("indexing", phase="catalog", drives_total=len(drives) * 2, drives_done=0, files_done=0)

    for i, drive_name in enumerate(drives):
        _prog.update("indexing", current_drive=drive_name, drives_done=i)

        sync_row = db.execute(
            select(DriveSync).where(DriveSync.drive_name == drive_name)
        ).scalar_one_or_none()
        stored_delta = sync_row.catalog_delta_link if sync_row else None

        try:
            try:
                changed_files, deleted_ids, new_delta_link = await list_files_delta(
                    **sp_kwargs, drive_name=drive_name, delta_link=stored_delta
                )
            except DeltaExpiredError:
                logger.warning(f"[catalog:{drive_name}] Delta token expired — full crawl")
                changed_files, deleted_ids, new_delta_link = await list_files_delta(
                    **sp_kwargs, drive_name=drive_name, delta_link=None
                )
        except Exception as e:
            logger.error(f"[catalog:{drive_name}] crawl failed: {e}")
            _prog.update("indexing", drives_done=i + 1)
            continue

        if deleted_ids:
            db.execute(delete(FileCatalog).where(FileCatalog.item_id.in_(deleted_ids)))

        for f in changed_files:
            item_id = f.get("id")
            if not item_id:
                continue
            name = f.get("name", "")
            folder_path = f.get("folder_path", "")
            ext = ("." + name.rsplit(".", 1)[-1].lower())[:50] if "." in name else ""
            norm = _normalize(f"{drive_name}/{folder_path}/{name}")

            modified = None
            mod_str = f.get("modified", "")
            if mod_str:
                try:
                    modified = datetime.fromisoformat(mod_str.replace("Z", "+00:00")).replace(tzinfo=None)
                except Exception:
                    modified = None

            existing = db.execute(
                select(FileCatalog).where(FileCatalog.item_id == item_id)
            ).scalar_one_or_none()
            if existing:
                existing.drive = drive_name
                existing.folder_path = folder_path
                existing.name = name
                existing.name_norm = norm
                existing.ext = ext
                existing.web_url = f.get("web_url")
                existing.size = f.get("size") or 0
                existing.modified = modified
                existing.indexed_at = now
            else:
                db.add(FileCatalog(
                    item_id=item_id,
                    drive=drive_name,
                    folder_path=folder_path,
                    name=name,
                    name_norm=norm,
                    ext=ext,
                    web_url=f.get("web_url"),
                    size=f.get("size") or 0,
                    modified=modified,
                    indexed_at=now,
                ))
            total += 1

        if new_delta_link:
            if sync_row is None:
                sync_row = DriveSync(drive_name=drive_name)
                db.add(sync_row)
            sync_row.catalog_delta_link = new_delta_link
            sync_row.last_catalog_sync = now.replace(tzinfo=None)

        db.commit()
        db.expunge_all()  # clear session identity map — prevents memory accumulation across drives
        logger.info(f"[catalog:{drive_name}] {len(changed_files)} changed, {len(deleted_ids)} deleted")
        _prog.update("indexing", drives_done=i + 1, files_done=total)

    logger.info(f"File catalog: {total} files across {len(drives)} drives")
    return total


async def _index_sharepoint(
    db: Session,
    archived_item_ids: set[str] | None = None,
    archived_paths_by_drive: dict[str, set[str]] | None = None,
    no_index_drives: set[str] | None = None,
    no_index_item_ids: set[str] | None = None,
    no_index_paths_by_drive: dict[str, set[str]] | None = None,
) -> int:
    if not all([
        settings.sharepoint_tenant_id,
        settings.sharepoint_client_id,
        settings.sharepoint_client_secret,
    ]):
        return 0

    import httpx

    drives_raw = settings.sharepoint_index_drives
    if drives_raw:
        drives = [d.strip() for d in drives_raw.split(",") if d.strip()]
        logger.info(f"SHAREPOINT_INDEX_DRIVES set — content-indexing {len(drives)} drives")
    else:
        drives = _CONTENT_DRIVES_DEFAULT
        logger.info(f"SHAREPOINT_INDEX_DRIVES not set — using default knowledge drives ({len(drives)} drives)")

    if not drives:
        logger.warning("No drives found for indexing")
        return 0

    sp_kwargs = dict(
        tenant_id=settings.sharepoint_tenant_id,
        client_id=settings.sharepoint_client_id,
        client_secret=settings.sharepoint_client_secret,
        site_hostname=settings.sharepoint_site_hostname,
        site_path=settings.sharepoint_site_path,
    )

    now = datetime.now(timezone.utc)
    total = 0

    from app import progress as _prog
    from app.models import DriveSync
    from app.sharepoint import list_files_delta, DeltaExpiredError
    n_drives = len(drives)
    # drives_total was set to len(drives)*2 by _build_file_catalog; content phase starts at halfway
    catalog_offset = n_drives

    async with httpx.AsyncClient(timeout=60) as http:
        for i, drive_name in enumerate(drives):
            logger.info(f"Indexing drive: {drive_name}")
            _prog.update("indexing", phase="content", current_drive=drive_name, drives_done=catalog_offset + i)

            if no_index_drives and drive_name in no_index_drives:
                logger.info(f"[{drive_name}] Skipped — excluded from indexing")
                _prog.update("indexing", drives_done=catalog_offset + i + 1)
                continue

            try:
                sync_row = db.execute(
                    select(DriveSync).where(DriveSync.drive_name == drive_name)
                ).scalar_one_or_none()
                stored_delta = sync_row.content_delta_link if sync_row else None

                try:
                    changed_files, deleted_ids, new_delta_link = await list_files_delta(
                        **sp_kwargs, drive_name=drive_name, delta_link=stored_delta
                    )
                except DeltaExpiredError:
                    logger.warning(f"[{drive_name}] Content delta expired — full crawl")
                    changed_files, deleted_ids, new_delta_link = await list_files_delta(
                        **sp_kwargs, drive_name=drive_name, delta_link=None
                    )

                count = await _index_one_drive(
                    db, drive_name, sp_kwargs, http, now,
                    changed_items=changed_files, deleted_item_ids=deleted_ids,
                    archived_item_ids=archived_item_ids,
                    archived_paths=archived_paths_by_drive.get(drive_name) if archived_paths_by_drive else None,
                    no_index_item_ids=no_index_item_ids,
                    no_index_paths=no_index_paths_by_drive.get(drive_name) if no_index_paths_by_drive else None,
                )
                logger.info(f"[{drive_name}] Indexed {count} chunks")
                total += count

                if new_delta_link:
                    if sync_row is None:
                        sync_row = db.execute(
                            select(DriveSync).where(DriveSync.drive_name == drive_name)
                        ).scalar_one_or_none()
                    if sync_row is None:
                        sync_row = DriveSync(drive_name=drive_name)
                        db.add(sync_row)
                    sync_row.content_delta_link = new_delta_link
                    sync_row.last_content_sync = now.replace(tzinfo=None)
                    db.commit()
                    db.expunge_all()

            except Exception as e:
                logger.error(f"[{drive_name}] Drive indexing failed: {e}")
            _prog.update("indexing", drives_done=catalog_offset + i + 1, files_done=total)

    return total


async def _index_doc_comments(db: Session) -> int:
    from app.models import DocMeta

    metas = db.execute(
        select(DocMeta).where(DocMeta.comment.isnot(None)).where(DocMeta.no_index == False)
    ).scalars().all()

    now = datetime.now(timezone.utc)
    indexed = 0

    for meta in metas:
        if not meta.comment:
            continue

        display = f"{meta.drive}/{meta.path}" if meta.path else meta.name
        content = f"[Document: {display}]\n\nAnnotation: {meta.comment}"

        try:
            embedding = await embed_text(content)
        except Exception as e:
            logger.error(f"Embedding failed for doc comment {meta.item_id}: {e}")
            continue

        existing = db.execute(
            select(KnowledgeChunk).where(
                KnowledgeChunk.source_type == "doc_comment",
                KnowledgeChunk.source_id == meta.item_id,
            )
        ).scalar_one_or_none()

        if existing:
            existing.content = content
            existing.source_name = display
            existing.embedding = embedding
            existing.indexed_at = now
        else:
            db.add(KnowledgeChunk(
                source_type="doc_comment",
                source_id=meta.item_id,
                source_name=display,
                source_url=None,
                content=content,
                embedding=embedding,
                indexed_at=now,
            ))
        indexed += 1

    db.commit()
    return indexed


async def _run_indexing() -> dict:
    if not settings.openai_api_key:
        return {"ok": False, "message": "OPENAI_API_KEY not set — skipping indexing."}

    from app import progress as _prog
    _prog.update("indexing", running=True, drives_done=0, drives_total=0, files_done=0, current_drive="", phase="")

    db = SessionLocal()
    try:
        # Remove any previously indexed Invenias chunks — Ask PL is SharePoint-only
        db.execute(delete(KnowledgeChunk).where(KnowledgeChunk.source_type == "invenias"))
        db.commit()

        # Load archive + no_index DocMeta state once before indexing begins
        from app.models import DocMeta as _DocMeta
        all_meta_rows = db.execute(select(_DocMeta)).scalars().all()

        archived_item_ids: set[str] = set()
        archived_paths_by_drive: dict[str, set[str]] = {}
        no_index_drives: set[str] = set()
        no_index_item_ids: set[str] = set()
        no_index_paths_by_drive: dict[str, set[str]] = {}

        for r in all_meta_rows:
            if r.is_archive:
                archived_item_ids.add(r.item_id)
                if r.drive and r.path:
                    archived_paths_by_drive.setdefault(r.drive, set()).add(r.path)
            if r.no_index:
                drive_name = _drive_for_meta_key(r.item_id)
                if drive_name:
                    no_index_drives.add(drive_name)
                else:
                    no_index_item_ids.add(r.item_id)
                    if r.drive and r.path:
                        no_index_paths_by_drive.setdefault(r.drive, set()).add(r.path)

        import gc
        catalog_count = await _build_file_catalog(db)
        gc.collect()
        sp_count = await _index_sharepoint(
            db, archived_item_ids, archived_paths_by_drive,
            no_index_drives, no_index_item_ids, no_index_paths_by_drive,
        )
        gc.collect()
        comment_count = await _index_doc_comments(db)
        msg = (
            f"{catalog_count} files cataloged, {sp_count} SharePoint chunks indexed, "
            f"{comment_count} doc comments indexed"
        )
        logger.info(msg)
        _prog.update("indexing", running=False, current_drive="",
                     last_finished_at=datetime.now(timezone.utc).isoformat(),
                     last_message=msg, last_ok=True)
        return {"ok": True, "message": msg}
    except Exception as e:
        db.rollback()
        logger.error(f"Indexing failed: {e}")
        _prog.update("indexing", running=False, current_drive="",
                     last_finished_at=datetime.now(timezone.utc).isoformat(),
                     last_message=str(e), last_ok=False)
        return {"ok": False, "message": str(e)}
    finally:
        db.close()


def run_indexing_sync() -> dict:
    return asyncio.run(_run_indexing())
